from __future__ import annotations

from pathlib import Path
import re

from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt


def normalize_lines(text: str) -> list[str]:
    lines: list[str] = []
    seen: set[str] = set()
    for raw in text.splitlines():
        line = re.sub(r"\s+", " ", raw).strip()
        if not line:
            continue
        if line in seen:
            continue
        seen.add(line)
        lines.append(line)
    return lines


def extract_slide_texts(html: str) -> list[tuple[str, list[str]]]:
    soup = BeautifulSoup(html, "html.parser")
    deck = soup.find("div", id="deck")
    if deck is None:
        raise ValueError("Could not find deck container with id='deck'.")

    extracted: list[tuple[str, list[str]]] = []
    for idx, slide in enumerate(deck.find_all("div", class_="slide"), start=1):
        title_tag = slide.find(["h1", "h2", "h3"])
        title = title_tag.get_text(" ", strip=True) if title_tag else f"Slide {idx}"

        # Remove title so body does not duplicate it.
        if title_tag is not None:
            title_tag.extract()

        body_text = slide.get_text("\n", strip=True)
        lines = normalize_lines(body_text)
        extracted.append((title, lines))

    return extracted


def add_body_text(placeholder, lines: list[str]) -> None:
    tf = placeholder.text_frame
    tf.clear()

    if not lines:
        return

    # Keep content readable on a 16:9 slide.
    max_lines = 18
    clipped = lines[:max_lines]

    for i, line in enumerate(clipped):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(18 if i < 3 else 16)


def html_to_pptx(input_html: Path, output_pptx: Path) -> None:
    html = input_html.read_text(encoding="utf-8")
    slides_data = extract_slide_texts(html)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    layout = prs.slide_layouts[1]  # Title and Content
    for title, lines in slides_data:
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        add_body_text(slide.placeholders[1], lines)

    prs.save(str(output_pptx))


if __name__ == "__main__":
    root = Path(__file__).resolve().parent
    input_file = root / "presentation.html"
    output_file = root / "InsureTrust_Presentation.pptx"

    html_to_pptx(input_file, output_file)
    print(f"Created: {output_file}")
